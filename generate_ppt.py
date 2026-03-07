"""
Petpooja AI Copilot — HACKaMINeD Hackathon Presentation Generator
Generates a polished 10-slide PowerPoint presentation.
Run: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, textwrap

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x0F, 0x1A)   # near-black navy
ACCENT     = RGBColor(0xF9, 0x74, 0x15)   # Petpooja orange
ACCENT2    = RGBColor(0x6C, 0x63, 0xFF)   # soft purple
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xDD)
MID_GREY   = RGBColor(0x44, 0x44, 0x66)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
CARD_BG    = RGBColor(0x1A, 0x1A, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

HACKAMINED_LOGO = os.path.join("docs", "hackamined_logo.png")   # saved by user
NIRMA_LOGO      = "nirma_logo.png"
BU_LOGO         = "bu_logo.png"
PETPOOJA_LOGO   = os.path.join("frontend", "public", "petpoja.png")

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h, bold=False, italic=False,
                size=Pt(14), color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = size
    run.font.color.rgb = color
    return txb


def add_logo(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)


def dark_slide(title_text, subtitle_text=None):
    """Creates a base dark slide with an orange top bar and optional subtitle."""
    slide = prs.slides.add_slide(BLANK)

    # background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)

    # top accent bar (full width, 6px)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=ACCENT)

    # bottom accent bar
    add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), fill=ACCENT2)

    if title_text:
        add_textbox(slide, title_text,
                    Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                    bold=True, size=Pt(28), color=ACCENT, align=PP_ALIGN.LEFT)
        # underline bar beneath title
        add_rect(slide, Inches(0.6), Inches(0.95), Inches(2.5), Inches(0.04), fill=ACCENT)

    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.6), Inches(1.0), Inches(12), Inches(0.45),
                    size=Pt(13), color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    return slide


def bullet_box(slide, items, l, t, w, h, bullet="▸", size=Pt(14), gap=Inches(0.38)):
    """Renders a list of bullet strings inside a textbox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size  = size
        run.font.color.rgb = LIGHT_GREY


def card(slide, l, t, w, h, title, body_lines, title_color=ACCENT):
    """Draws a dark card with a coloured heading and body lines."""
    add_rect(slide, l, t, w, h, fill=CARD_BG, line=MID_GREY, line_w=Pt(0.75))
    add_textbox(slide, title,
                l + Inches(0.18), t + Inches(0.12), w - Inches(0.3), Inches(0.38),
                bold=True, size=Pt(13), color=title_color)
    bullet_box(slide, body_lines,
               l + Inches(0.18), t + Inches(0.5), w - Inches(0.3), h - Inches(0.55),
               size=Pt(11.5))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)

# top gradient bar
add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=ACCENT)
# bottom bar
add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), fill=ACCENT2)

# ── Logos row ────────────────────────────────────────────────────────────
logo_y = Inches(0.22)
logo_h = Inches(0.88)

# HACKaMINeD logo (provided) — centre
add_logo(slide, HACKAMINED_LOGO, SLIDE_W/2 - Inches(0.6), logo_y, Inches(1.2), logo_h)
# Nirma University logo (left)
add_logo(slide, NIRMA_LOGO, Inches(0.4), logo_y, Inches(1.5), logo_h)
# Binghamton University logo (right)
add_logo(slide, BU_LOGO, SLIDE_W - Inches(2.2), logo_y, Inches(1.7), logo_h)
# Petpooja logo (track sponsor)
add_logo(slide, PETPOOJA_LOGO, SLIDE_W - Inches(4.2), logo_y, Inches(1.5), logo_h)

# ── Divider ──────────────────────────────────────────────────────────────
add_rect(slide, Inches(0.5), Inches(1.28), SLIDE_W - Inches(1.0), Inches(0.025), fill=MID_GREY)

# ── Institute + Track badge ───────────────────────────────────────────────
add_textbox(slide, "Nirma University  •  Institute of Technology",
            Inches(0.5), Inches(1.35), Inches(8), Inches(0.38),
            size=Pt(12.5), color=LIGHT_GREY, align=PP_ALIGN.LEFT)

# Track badge pill (orange rounded rect)
add_rect(slide, Inches(8.9), Inches(1.32), Inches(4.0), Inches(0.36), fill=ACCENT)
add_textbox(slide, "Track: AI-Powered Revenue & Voice Copilot",
            Inches(8.92), Inches(1.33), Inches(3.96), Inches(0.34),
            bold=True, size=Pt(10.5), color=WHITE, align=PP_ALIGN.CENTER)

# ── Track Sponsor line ────────────────────────────────────────────────────
add_textbox(slide, "Track Sponsor: Petpooja",
            Inches(0.5), Inches(1.78), Inches(6), Inches(0.32),
            size=Pt(11.5), color=ACCENT, align=PP_ALIGN.LEFT)

# ── Main Title ───────────────────────────────────────────────────────────
add_textbox(slide, "🍽️  Petpooja AI Copilot",
            Inches(0.5), Inches(2.25), Inches(12.3), Inches(1.1),
            bold=True, size=Pt(46), color=WHITE, align=PP_ALIGN.CENTER)

# orange underline
add_rect(slide, Inches(3.5), Inches(3.35), Inches(6.3), Inches(0.06), fill=ACCENT)

# ── Tagline ───────────────────────────────────────────────────────────────
add_textbox(slide, "AI-Powered Revenue Intelligence · Voice Ordering · Inventory Automation",
            Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
            size=Pt(16), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ── Hackathon name ────────────────────────────────────────────────────────
add_textbox(slide, "HACKaMINeD 2026",
            Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.42),
            bold=True, size=Pt(18), color=ACCENT2, align=PP_ALIGN.CENTER)

# ── Team box ─────────────────────────────────────────────────────────────
add_rect(slide, Inches(3.0), Inches(4.62), Inches(7.3), Inches(1.65), fill=CARD_BG, line=ACCENT, line_w=Pt(1.2))
add_textbox(slide, "Team Name:  [Your Team Name]",
            Inches(3.2), Inches(4.72), Inches(7.0), Inches(0.4),
            bold=True, size=Pt(14), color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slide,
            "Member 1  ·  Member 2  ·  Member 3  ·  Member 4",
            Inches(3.2), Inches(5.18), Inches(7.0), Inches(0.38),
            size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Nirma University, Ahmedabad",
            Inches(3.2), Inches(5.6), Inches(7.0), Inches(0.35),
            size=Pt(11.5), color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Problem Statement",
                   "The Hidden Crisis in India's Restaurant Industry")

# large stat cards row
stats = [
    ("~₹40,000 Cr",  "Lost annually due to\nunoptimised menu pricing"),
    ("< 20%",        "Restaurants use any\ndata-driven analytics tool"),
    ("2–3 hrs/day",  "Wasted on manual\nphone order taking"),
    ("30–40%",       "Food waste from poor\ninventory tracking"),
]
cw = Inches(2.9)
cy = Inches(1.6)
for i, (val, lbl) in enumerate(stats):
    cx = Inches(0.45) + i * (cw + Inches(0.25))
    add_rect(slide, cx, cy, cw, Inches(1.7), fill=CARD_BG, line=ACCENT, line_w=Pt(1.2))
    add_textbox(slide, val, cx + Inches(0.1), cy + Inches(0.18), cw - Inches(0.2), Inches(0.65),
                bold=True, size=Pt(28), color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, lbl, cx + Inches(0.1), cy + Inches(0.82), cw - Inches(0.2), Inches(0.7),
                size=Pt(11), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# narrative block
prob_text = (
    "Restaurant owners in India manage 40–60 menu items and hundreds of daily transactions "
    "with zero visibility into which items are profitable, which are dragging margins down, "
    "and which hidden gems are being ignored. Phone-based ordering — still dominant in India — "
    "is slow, error-prone, and language-mixed (Hinglish). Inventory is managed on paper or basic "
    "spreadsheets, leading to stockouts, over-ordering, and silent revenue leakage.\n\n"
    "There is NO affordable, locally-running AI solution designed specifically for Indian SMB restaurants."
)
add_rect(slide, Inches(0.45), Inches(3.55), Inches(12.4), Inches(2.7), fill=CARD_BG, line=MID_GREY, line_w=Pt(0.6))
add_textbox(slide, prob_text,
            Inches(0.65), Inches(3.7), Inches(12.0), Inches(2.5),
            size=Pt(13), color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Introduction",
                   "What is Petpooja AI Copilot?")

intro_blurb = (
    "Petpooja AI Copilot is a full-stack, locally-running AI platform that transforms raw "
    "Point-of-Sale data into intelligent revenue decisions — and automates Hinglish voice "
    "ordering — with zero cloud API costs."
)
add_textbox(slide, intro_blurb,
            Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.75),
            size=Pt(14), color=LIGHT_GREY)

# 6 module cards
modules = [
    ("📊 Revenue Intelligence", "Menu Matrix (Star/Hidden Star/\nPlowhorse/Dog) · Margin analysis\n· Price optimisation"),
    ("🎙️ Voice Ordering Copilot", "Whisper STT (local) · Hinglish\nfuzzy matching · Upsell engine\n· KOT generation"),
    ("📦 Inventory Manager", "Real-time stock tracking · Low-stock\nalerts · Supplier management\n· Auto-depletion on order"),
    ("🧾 Recipe / BOM", "Bill of Materials per item·\nAuto ingredient deduction\n· 51 pre-seeded recipes"),
    ("🍽️ Customer Menu + KDS", "Filterable menu · Cart & GST\ncheckout · Kitchen Display\nSystem with order queue"),
    ("🧩 Combo Engine", "Apriori association mining\non 800 transactions ·\nAOV-boosting bundle suggestions"),
]
cols, rows = 3, 2
cw, ch = Inches(4.1), Inches(1.75)
for i, (ttl, body) in enumerate(modules):
    col = i % cols
    row = i // cols
    cx = Inches(0.45) + col * (cw + Inches(0.2))
    cy = Inches(2.05) + row * (ch + Inches(0.18))
    card(slide, cx, cy, cw, ch, ttl, body.split("\n"))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Proposed Approach
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Proposed Approach",
                   "Architecture & Algorithm Pipeline")

# architecture diagram (text-based)
arch_lines = [
    "┌──────────────── React Frontend  (Vite · Tailwind CSS · Recharts) ────────────────┐",
    "│  Dashboard  │  Menu Intelligence  │  Combo Engine  │  Voice Orders  │  Inventory │",
    "└───────────────────────────────────┬───────────────────────────────────────────────┘",
    "                                    │  REST API  (Vite proxy → localhost:8000)       ",
    "┌───────────────────────────────────┴───────────────────────────────────────────────┐",
    "│                         FastAPI + Uvicorn  Backend                                │",
    "│  ┌──Revenue Engine──────────────┐   ┌──Voice Copilot────────────────────────────┐ │",
    "│  │ • Pandas margin & velocity   │   │ • OpenAI Whisper STT (base, local)        │ │",
    "│  │ • Menu Matrix classification │   │ • RapidFuzz fuzzy match (≥65% threshold)  │ │",
    "│  │ • Apriori mining (mlxtend)   │   │ • Hinglish qty parser (ek/do/teen)        │ │",
    "│  │ • Rule-based price optimizer │   │ • Upsell engine + KOT generator           │ │",
    "│  └──────────────────────────────┘   └───────────────────────────────────────────┘ │",
    "│                     JSON Data Store  (menu · inventory · recipes · orders)        │",
    "└───────────────────────────────────────────────────────────────────────────────────┘",
]
add_rect(slide, Inches(0.35), Inches(1.1), Inches(12.6), Inches(4.05), fill=RGBColor(0x0D,0x0D,0x1C), line=ACCENT2, line_w=Pt(0.8))
txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.35), Inches(4.0))
tf = txb.text_frame
tf.word_wrap = False
first = True
for line in arch_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    run = p.add_run()
    run.text = line
    run.font.size = Pt(9.5)
    run.font.color.rgb = GREEN
    run.font.name = "Courier New"

# tech stack pills row
stack = ["Python 3.12", "FastAPI", "OpenAI Whisper", "RapidFuzz", "MLxtend / Apriori",
         "Pandas · NumPy", "React 19 · Vite 5", "Tailwind CSS 4", "Recharts · Framer Motion"]
pill_x = Inches(0.35)
pill_y = Inches(5.3)
for tech in stack:
    pw = Inches(len(tech) * 0.105 + 0.25)
    add_rect(slide, pill_x, pill_y, pw, Inches(0.34), fill=ACCENT2)
    add_textbox(slide, tech, pill_x + Inches(0.08), pill_y + Inches(0.03), pw - Inches(0.1), Inches(0.3),
                bold=True, size=Pt(9.5), color=WHITE, align=PP_ALIGN.CENTER)
    pill_x += pw + Inches(0.1)
    if pill_x > Inches(12.5):
        pill_x = Inches(0.35)
        pill_y += Inches(0.42)

add_textbox(slide, "⚡  100% Local — No external LLM API calls · No cloud costs",
            Inches(0.5), Inches(6.55), Inches(12.0), Inches(0.38),
            bold=True, size=Pt(13), color=ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Why Choose Our Solution?
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Why Choose Petpooja AI Copilot?",
                   "Key differentiators over existing restaurant management tools")

usp_cards = [
    ("🆓  Zero API Cost",
     ["Runs 100% locally — Whisper, Apriori,", "fuzzy matching all on-device",
      "No monthly subscription to OpenAI / cloud"]),
    ("🗣️  Hinglish-Native Voice AI",
     ["Handles real Indian restaurant orders:", "\"do paneer burger aur ek large fries dena\"",
      "Quantity words: ek/do/teen/chaar mapped"]),
    ("💎  Hidden Star Discovery",
     ["Finds high-margin, under-promoted items", "Tells owner exactly which dishes to push",
      "Direct actionable revenue uplift"]),
    ("📦  Auto Inventory Depletion",
     ["Recipe BOM linked to every menu item", "Stock auto-reduces on each confirmed order",
      "Low-stock alerts prevent stockouts"]),
    ("🧩  AI Combo Suggestions",
     ["Apriori mines 800 real transactions", "Bundles that statistically sell together",
      "Each combo shows AOV lift & margin score"]),
    ("🖥️  Full-Stack SaaS UX",
     ["Dark glassmorphism dashboard", "Kitchen Display System for staff",
      "Customer menu + cart + GST checkout"]),
]
cw, ch = Inches(4.0), Inches(1.85)
for i, (ttl, body) in enumerate(usp_cards):
    col = i % 3
    row = i // 3
    cx = Inches(0.4) + col * (cw + Inches(0.26))
    cy = Inches(1.45) + row * (ch + Inches(0.18))
    card(slide, cx, cy, cw, ch, ttl, body)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Limitations
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Limitations & Scope",
                   "Honest assessment of when the solution may face difficulties")

limitations = [
    ("🎤  Audio Quality", [
        "Whisper struggles with heavy background noise (busy kitchen, TV)",
        "Very strong regional accents or dialects beyond Hindi/English may reduce accuracy",
        "Optimal with a clear microphone from ~30 cm distance",
    ]),
    ("📊  Data Dependency", [
        "Revenue Intelligence requires at least ~200+ transactions for meaningful Apriori rules",
        "Menu Matrix classification uses median thresholds — small menus (<10 items) may skew results",
        "Current mock data is vegetarian-pizza-centric; real menus need re-seeding",
    ]),
    ("🔌  Hardware Requirements", [
        "Whisper 'base' model requires ~1 GB RAM; slower on CPUs without AVX2",
        "No mobile/tablet client — frontend is desktop-browser only in current version",
        "Requires Python 3.10+ and Node 18+ to be pre-installed locally",
    ]),
    ("🌐  Scope Boundaries", [
        "Single-outlet only — multi-branch consolidation not yet supported",
        "No payment gateway or real POS hardware integration in MVP",
        "Demand forecasting, collaborative filtering, anomaly detection are proposed but not implemented",
        "No real-time sync or offline-first PWA support",
    ]),
]
cw, ch = Inches(6.0), Inches(2.2)
for i, (ttl, body) in enumerate(limitations):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * (cw + Inches(0.5))
    cy = Inches(1.38) + row * (ch + Inches(0.22))
    card(slide, cx, cy, cw, ch, ttl, body)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Demo / Screenshots
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Screenshot & Demonstration",
                   "Live system running at http://localhost:5173")

screens = [
    ("📊 Revenue Dashboard",   Inches(0.4),  Inches(1.45), Inches(4.0), Inches(2.6)),
    ("🍽️ Menu Intelligence",   Inches(4.7),  Inches(1.45), Inches(4.0), Inches(2.6)),
    ("🎙️ Voice Ordering",      Inches(8.95), Inches(1.45), Inches(4.0), Inches(2.6)),
    ("📦 Inventory Manager",   Inches(0.4),  Inches(4.3),  Inches(4.0), Inches(2.6)),
    ("🧩 Combo Engine",        Inches(4.7),  Inches(4.3),  Inches(4.0), Inches(2.6)),
    ("🖥️ Kitchen Display",     Inches(8.95), Inches(4.3),  Inches(4.0), Inches(2.6)),
]
for (lbl, x, y, w, h) in screens:
    add_rect(slide, x, y, w, h, fill=CARD_BG, line=MID_GREY, line_w=Pt(0.75))
    add_textbox(slide, lbl, x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.35),
                bold=True, size=Pt(12), color=ACCENT)
    add_textbox(slide, "[ Screenshot / Screen recording here ]",
                x + Inches(0.1), y + Inches(0.55), w - Inches(0.2), h - Inches(0.7),
                size=Pt(11), color=MID_GREY, align=PP_ALIGN.CENTER)

add_textbox(slide,
            "▶  Live demo: python -m uvicorn main:app --port 8000  +  npm run dev  →  localhost:5173",
            Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.32),
            size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ML Models (Additional Info)
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("ML Models — Implemented & Proposed",
                   "Currently running + roadmap for future intelligence")

# implemented table
impl = [
    ("OpenAI Whisper",      "Speech-to-text from audio",               "whisper",    "✅ Live"),
    ("RapidFuzz",           "Fuzzy menu item matching (partial ratio)", "rapidfuzz",  "✅ Live"),
    ("Apriori Algorithm",   "Combo/bundle recommendations",            "mlxtend",    "✅ Live"),
    ("Menu Matrix",         "Star / Hidden Star / Plowhorse / Dog",    "pandas",     "✅ Live"),
]
proposed = [
    ("Prophet / SARIMA",    "Demand forecasting for inventory",        "prophet",    "🔜 Proposed"),
    ("NLP Intent Parser",   "LLM-based complex order understanding",   "transformers","🔜 Proposed"),
    ("Collab Filtering",    "Personalised upsell recommendations",     "scikit-surprise","🔜 Proposed"),
    ("Anomaly Detection",   "Revenue / inventory anomaly alerts",      "scikit-learn", "🔜 Proposed"),
]
headers = ["Model", "Purpose", "Library", "Status"]
col_ws  = [Inches(2.4), Inches(4.5), Inches(2.5), Inches(1.8)]
row_h   = Inches(0.36)
table_x = Inches(0.4)

def draw_table(slide, rows_data, start_y, section_label, label_color):
    add_textbox(slide, section_label,
                table_x, start_y - Inches(0.32), Inches(4), Inches(0.3),
                bold=True, size=Pt(12), color=label_color)
    # header
    hx = table_x
    for hi, hdr in enumerate(headers):
        add_rect(slide, hx, start_y, col_ws[hi], row_h, fill=MID_GREY)
        add_textbox(slide, hdr, hx + Inches(0.06), start_y + Inches(0.04),
                    col_ws[hi] - Inches(0.1), row_h - Inches(0.08),
                    bold=True, size=Pt(11), color=WHITE)
        hx += col_ws[hi]
    # data rows
    for ri, row in enumerate(rows_data):
        ry = start_y + (ri + 1) * row_h
        rx = table_x
        bg = CARD_BG if ri % 2 == 0 else RGBColor(0x12, 0x12, 0x24)
        for ci, cell in enumerate(row):
            add_rect(slide, rx, ry, col_ws[ci], row_h, fill=bg, line=MID_GREY, line_w=Pt(0.4))
            col_color = GREEN if "✅" in cell else (ACCENT2 if "🔜" in cell else LIGHT_GREY)
            add_textbox(slide, cell, rx + Inches(0.06), ry + Inches(0.04),
                        col_ws[ci] - Inches(0.1), row_h - Inches(0.08),
                        size=Pt(10.5), color=col_color)
            rx += col_ws[ci]

draw_table(slide, impl,     Inches(1.32), "Currently Implemented", GREEN)
draw_table(slide, proposed, Inches(4.15), "Proposed / Roadmap", ACCENT2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture & Tech Deep-Dive
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Technical Deep-Dive",
                   "System design decisions and data flow")

deep_cards = [
    ("⚙️  Revenue Engine Pipeline", [
        "1. Load menu.json + transactions.json via db.py",
        "2. Pandas: compute contribution_margin() per item",
        "3. Classify velocity (Fast / Moderate / Slow) by percentile",
        "4. Apply 2×2 Menu Matrix using median thresholds",
        "5. Run Apriori (min_support=0.01, min_confidence=0.30)",
        "6. Overlay Hidden Star flag on combo suggestions",
    ]),
    ("🎙️  Voice Order Pipeline", [
        "1. Audio uploaded → POST /voice/transcribe",
        "2. Whisper base model → raw text (multilingual)",
        "3. Hinglish parser: ek→1, do→2, teen→3 ...",
        "4. RapidFuzz partial_ratio (≥65%) → menu item match",
        "5. Modifier extraction: 'extra cheese', 'large', 'spicy'",
        "6. Upsell check → POST /voice/confirm-order → KOT",
    ]),
    ("📦  Inventory Auto-Depletion", [
        "Each menu item has a Bill of Materials in recipes.json",
        "On order confirmation: for each item × qty → BOM lookup",
        "Subtract ingredient quantities from inventory.json",
        "If stock < min_threshold → append to low_stock_alerts[]",
        "Real-time alerts surfaced in Inventory Manager UI",
    ]),
    ("🗃️  Data Store Strategy", [
        "JSON file store chosen for zero-setup portability",
        "db.py provides a clean read/write abstraction layer",
        "Easily swappable to SQLite / PostgreSQL for production",
        "51 menu items · 20 ingredients · 800 transactions pre-seeded",
    ]),
]
cw, ch = Inches(6.0), Inches(2.72)
for i, (ttl, body) in enumerate(deep_cards):
    col = i % 2
    row = i // 2
    cx = Inches(0.4) + col * (cw + Inches(0.5))
    cy = Inches(1.38) + row * (ch + Inches(0.22))
    card(slide, cx, cy, cw, ch, ttl, body)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — References
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("References",
                   "Libraries, papers, and resources used in this project")

refs = [
    ("[1]  OpenAI Whisper",
     "Radford, A. et al. (2022). Robust Speech Recognition via Large-Scale Weak Supervision. arXiv:2212.04356.",
     "https://arxiv.org/abs/2212.04356"),
    ("[2]  MLxtend Apriori",
     "Raschka, S. (2018). MLxtend: Providing machine learning and data science utilities and extensions. JOSS.",
     "https://rasbt.github.io/mlxtend/"),
    ("[3]  RapidFuzz",
     "Bachthaler, M. (2020). RapidFuzz: Rapid fuzzy string matching in Python using various string metrics.",
     "https://github.com/maxbachmann/RapidFuzz"),
    ("[4]  Menu Engineering",
     "Kasavana, M. L. & Smith, D. I. (1982). Menu Engineering: A Practical Guide to Menu Analysis.",
     "Hospitality Publications, Michigan State University"),
    ("[5]  FastAPI",
     "Ramírez, S. (2018). FastAPI — Modern, fast (high-performance) web framework for building APIs with Python.",
     "https://fastapi.tiangolo.com"),
    ("[6]  React / Vite",
     "Meta Open Source (2013–2024). React — A JavaScript library for building user interfaces.",
     "https://react.dev  ·  https://vitejs.dev"),
    ("[7]  Petpooja Track Brief",
     "HACKaMINeD 2026 — AI-Powered Revenue & Voice Copilot for Restaurants track specification.",
     "Nirma University · Binghamton University"),
]
for i, (tag, text, url) in enumerate(refs):
    cy = Inches(1.38) + i * Inches(0.76)
    add_rect(slide, Inches(0.4), cy, Inches(12.5), Inches(0.7), fill=CARD_BG, line=MID_GREY, line_w=Pt(0.4))
    add_textbox(slide, tag, Inches(0.55), cy + Inches(0.05), Inches(1.8), Inches(0.3),
                bold=True, size=Pt(11), color=ACCENT)
    add_textbox(slide, text, Inches(2.4), cy + Inches(0.04), Inches(9.5), Inches(0.35),
                size=Pt(10.5), color=LIGHT_GREY)
    add_textbox(slide, url, Inches(2.4), cy + Inches(0.38), Inches(9.5), Inches(0.28),
                size=Pt(9.5), color=ACCENT2)


# ── Save ────────────────────────────────────────────────────────────────────
out_path = os.path.join("docs", "Petpooja_Presentation.pptx")
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
